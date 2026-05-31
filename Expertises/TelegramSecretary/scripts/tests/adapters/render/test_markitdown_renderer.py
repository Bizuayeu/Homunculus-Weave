"""MarkitdownRenderer の integration test。

実 markitdown ライブラリを呼んで md 化を検証する。fixture は python-docx /
openpyxl / python-pptx で test 内で動的生成（git に大きなバイナリを置かない）。
"""
from __future__ import annotations

import zipfile
from pathlib import Path

import pytest

from adapters.render.markitdown_renderer import MarkitdownRenderer
from domain.media import MediaAttachment


@pytest.fixture(scope="module")
def renderer() -> MarkitdownRenderer:
    """MarkItDown 初期化は magika ML model load が走り重いので module スコープ。"""
    return MarkitdownRenderer()


def _docx_media(file_id: str = "docx") -> MediaAttachment:
    return MediaAttachment(
        kind="document",
        file_id=file_id,
        mime_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        size=4096,
        file_name=f"{file_id}.docx",
    )


def _xlsx_media() -> MediaAttachment:
    return MediaAttachment(
        kind="document",
        file_id="xlsx",
        mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        size=4096,
        file_name="data.xlsx",
    )


def _pptx_media() -> MediaAttachment:
    return MediaAttachment(
        kind="document",
        file_id="pptx",
        mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        size=4096,
        file_name="slides.pptx",
    )


# === 正常系: 各 mime で md 化成功 ===


def test_renders_docx_to_markdown(tmp_path: Path, renderer: MarkitdownRenderer):
    """python-docx で fixture 作成 → markitdown で md 化、本文が含まれる。"""
    from docx import Document

    doc_path = tmp_path / "spec.docx"
    doc = Document()
    doc.add_heading("仕様書", level=1)
    doc.add_paragraph("これは render 検証用の docx です。")
    doc.add_paragraph("ブルーベリーは美味しい。")
    doc.save(str(doc_path))

    rendered = renderer.render(_docx_media(), doc_path)
    assert rendered.render_status == "ok"
    assert rendered.rendered_text is not None
    assert "仕様書" in rendered.rendered_text
    assert "ブルーベリー" in rendered.rendered_text


def test_renders_xlsx_to_markdown(tmp_path: Path, renderer: MarkitdownRenderer):
    """openpyxl で fixture 作成 → markitdown が表データを md 化。"""
    from openpyxl import Workbook

    xlsx_path = tmp_path / "data.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.append(["品種", "色"])
    ws.append(["Pink Popcorn", "ピンク"])
    ws.append(["Pearl", "白"])
    wb.save(str(xlsx_path))

    rendered = renderer.render(_xlsx_media(), xlsx_path)
    assert rendered.render_status == "ok"
    assert rendered.rendered_text is not None
    # markitdown は表をパイプ区切り or 類似形式に変換、品種名が含まれる
    assert "Pink Popcorn" in rendered.rendered_text
    assert "Pearl" in rendered.rendered_text


def test_renders_pptx_to_markdown(tmp_path: Path, renderer: MarkitdownRenderer):
    """python-pptx で fixture 作成 → markitdown がスライド内容を md 化。"""
    from pptx import Presentation

    pptx_path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    title_box = slide.shapes.add_textbox(left=0, top=0, width=5000000, height=500000)
    title_box.text_frame.text = "サンプルプロジェクト"
    prs.save(str(pptx_path))

    rendered = renderer.render(_pptx_media(), pptx_path)
    assert rendered.render_status == "ok"
    assert rendered.rendered_text is not None
    assert "サンプル" in rendered.rendered_text


# === 失敗系: クラッシュさせず flag 化 ===


def test_renders_returns_failed_for_empty_docx(
    tmp_path: Path, renderer: MarkitdownRenderer
):
    """空ファイル（.docx 拡張子）→ markitdown 内部で BadZipFile → render_status='failed'。"""
    empty = tmp_path / "empty.docx"
    empty.write_bytes(b"")

    rendered = renderer.render(_docx_media("empty"), empty)
    assert rendered.render_status == "failed"
    assert rendered.rendered_text is None


def test_renders_returns_failed_for_nonexistent_file(
    tmp_path: Path, renderer: MarkitdownRenderer
):
    """存在しないファイル → markitdown が FileNotFoundError → render_status='failed'。"""
    missing = tmp_path / "missing.docx"  # 作らない
    rendered = renderer.render(_docx_media("missing"), missing)
    assert rendered.render_status == "failed"
    assert rendered.rendered_text is None


def test_renders_garbage_bytes_returns_ok_as_text_passthrough(
    tmp_path: Path, renderer: MarkitdownRenderer
):
    """markitdown は寛容: garbage バイト列 (.docx 拡張子) でも text として何か返す。

    実挙動: magika ML model が file type を中身で推定し、plain text と判定すれば
    バイト列をそのまま rendered_text として返す。**この設計を エージェント 側に伝え**、
    エージェント が「意味のあるテキストか？」を判断する責務とする（LLM 推論をコード外に出す分業）。
    """
    garbage = tmp_path / "garbage.docx"
    garbage.write_bytes(b"this is not a valid docx, just random bytes")

    rendered = renderer.render(_docx_media("garbage"), garbage)
    assert rendered.render_status == "ok"
    assert rendered.rendered_text is not None
    assert "random bytes" in rendered.rendered_text


# === セキュリティ: 例外メッセージに絶対パス・file_id 全文を含めない ===


def test_warning_does_not_leak_absolute_path_or_full_file_id(
    tmp_path: Path, renderer: MarkitdownRenderer, capsys
):
    """render 失敗時の stderr warning に絶対パス・file_id 全文が出ない。"""
    # 空 .docx で BadZipFile を踏ませ、failed パスに入れる
    empty = tmp_path / "leak_check.docx"
    empty.write_bytes(b"")
    long_file_id = "AgACAgIAA_SENSITIVE_TAIL_DO_NOT_LEAK"

    renderer.render(_docx_media(long_file_id), empty)
    err = capsys.readouterr().err
    # 絶対パスは出さない
    assert str(empty) not in err
    assert "leak_check.docx" not in err
    # file_id は先頭 8 文字までで切れる
    assert "SENSITIVE_TAIL" not in err
    # ただし短縮 id（先頭 8 桁）は warning に残る（運用上の追跡性）
    assert long_file_id[:8] in err
